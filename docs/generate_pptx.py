"""Generate the M365 Agentic Service Developer Guide PowerPoint."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
BG_DARK = RGBColor(0x1B, 0x1B, 0x2F)
BG_CARD = RGBColor(0x26, 0x26, 0x42)
ACCENT = RGBColor(0x00, 0xB3, 0x36)
ACCENT2 = RGBColor(0x4F, 0xC3, 0xFF)
ACCENT3 = RGBColor(0xFF, 0xA7, 0x26)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xCC, 0xCC, 0xCC)
DIM = RGBColor(0x88, 0x88, 0xAA)
RED = RGBColor(0xFF, 0x66, 0x66)
PURPLE = RGBColor(0xBB, 0x66, 0xFF)
PINK = RGBColor(0xFF, 0x66, 0xFF)


def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def txt(slide, l, t, w, h, text, sz=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font="Segoe UI"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tb


def card(slide, l, t, w, h, title, lines, tc=ACCENT):
    rect(slide, l, t, w, h, BG_CARD, RGBColor(0x44, 0x44, 0x66))
    txt(slide, l + Inches(0.2), t + Inches(0.1), w - Inches(0.4), Inches(0.35), title, 14, tc, True)
    y = t + Inches(0.5)
    for line in lines:
        txt(slide, l + Inches(0.2), y, w - Inches(0.4), Inches(0.25), line, 11, LIGHT)
        y += Inches(0.26)


BLANK = 6

# ── SLIDE 1: Title ──────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[BLANK])
set_bg(s, BG_DARK)
rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)
txt(s, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
    "Deploying an Agentic Service\nto Microsoft 365 Copilot", 44, WHITE, True, PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(3.6), Inches(11), Inches(0.8),
    "with Delegated On-Behalf-Of (OBO) Access", 28, ACCENT, False, PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(5.0), Inches(11), Inches(0.6),
    "A general developer guideline for any agentic framework", 18, DIM, False, PP_ALIGN.CENTER)

# ── SLIDE 2: Architecture ───────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[BLANK])
set_bg(s, BG_DARK)
rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)
txt(s, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6), "Architecture Overview", 32, WHITE, True)
txt(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.4),
    "Two independently deployable services with clean trust boundaries", 16, DIM)

# M365 column
rect(s, Inches(0.5), Inches(1.7), Inches(3.0), Inches(2.4), RGBColor(0x1E, 0x3A, 0x5F), ACCENT2)
txt(s, Inches(0.7), Inches(1.8), Inches(2.6), Inches(0.3), "Microsoft 365", 14, ACCENT2, True)
card(s, Inches(0.7), Inches(2.2), Inches(2.6), Inches(0.7), "End User", ["Signs in via Entra SSO"], ACCENT2)
card(s, Inches(0.7), Inches(3.05), Inches(2.6), Inches(0.7), "M365 Copilot / Teams", ["Conversation UX + SSO"], ACCENT2)

# ACA column
rect(s, Inches(4.6), Inches(1.7), Inches(4.8), Inches(2.4), RGBColor(0x1A, 0x3D, 0x1A), ACCENT)
txt(s, Inches(4.8), Inches(1.8), Inches(4.4), Inches(0.3), "Azure Container Apps", 14, ACCENT, True)
card(s, Inches(4.8), Inches(2.2), Inches(2.1), Inches(0.7), "M365 Gateway", ["Stateless", "Bot protocol + JWT + OBO #1"], ACCENT)
card(s, Inches(7.1), Inches(2.2), Inches(2.1), Inches(0.7), "Agentic Service", ["Stateful, internal-only", "Your logic + OBO #2"], ACCENT)

# Downstream column
rect(s, Inches(10.2), Inches(1.7), Inches(2.6), Inches(2.4), RGBColor(0x4D, 0x33, 0x1A), ACCENT3)
txt(s, Inches(10.4), Inches(1.8), Inches(2.2), Inches(0.3), "Downstream", 14, ACCENT3, True)
card(s, Inches(10.4), Inches(2.2), Inches(2.2), Inches(0.7), "APIs / DBs / MCP", ["Called as delegated user"], ACCENT3)

# Why two services
rect(s, Inches(0.5), Inches(4.5), Inches(3.9), Inches(2.5), BG_CARD, RGBColor(0x44, 0x44, 0x66))
txt(s, Inches(0.7), Inches(4.6), Inches(3.5), Inches(0.35), "Why two services?", 18, ACCENT, True)
items = [
    ("Trust boundary separation", "Gateway handles Bot protocol. Service owns data access."),
    ("Independent scaling", "Gateway scales horizontally. Service pinned for sessions."),
    ("Framework freedom", "Any agentic framework inside the service."),
    ("Adapt without modifying", "Front an existing service by writing only the gateway."),
]
y = Inches(5.1)
for t, d in items:
    txt(s, Inches(0.7), y, Inches(3.5), Inches(0.25), t, 12, ACCENT2, True)
    txt(s, Inches(0.7), y + Inches(0.25), Inches(3.5), Inches(0.3), d, 11, LIGHT)
    y += Inches(0.55)

# Layer table
rect(s, Inches(4.8), Inches(4.5), Inches(7.9), Inches(2.5), BG_CARD, RGBColor(0x44, 0x44, 0x66))
txt(s, Inches(5.0), Inches(4.6), Inches(7.5), Inches(0.35), "Layer Responsibilities", 18, ACCENT, True)
layers = [
    ("M365 Copilot", "User identity, SSO, conversation UX", "None (platform)"),
    ("Gateway", "Bot protocol adapter, JWT validation, OBO #1", "Stateless"),
    ("Agentic Service", "Business logic, orchestration, sessions", "Stateful"),
    ("Downstream", "Data, APIs, MCP servers as delegated user", "External"),
]
y = Inches(5.1)
for layer, resp, state in layers:
    txt(s, Inches(5.0), y, Inches(2.0), Inches(0.3), layer, 12, ACCENT2, True)
    txt(s, Inches(7.0), y, Inches(3.8), Inches(0.3), resp, 12, LIGHT)
    txt(s, Inches(10.8), y, Inches(1.6), Inches(0.3), state, 12, DIM)
    y += Inches(0.38)

# ── SLIDE 3: Token Flow ─────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[BLANK])
set_bg(s, BG_DARK)
rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)
txt(s, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6), "Token Flow Deep Dive", 32, WHITE, True)
txt(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.4),
    "User identity flows end-to-end through two chained OBO exchanges", 16, DIM)

boxes = [
    ("End User", "Signs in\n(Entra SSO)", ACCENT2, 0.3),
    ("Copilot", "Activity +\nchannel token", ACCENT2, 2.7),
    ("Gateway", "OBO #1 + JWT validation\nextract claims\nforward token + headers", ACCENT, 5.1),
    ("Service", "Trusts gateway\nBinds ContextVar\nOBO #2 \u2192 downstream", ACCENT, 7.5),
    ("Downstream", "API call as\ndelegated user", ACCENT3, 10.3),
]
for title, body, color, left in boxes:
    rect(s, Inches(left), Inches(1.8), Inches(2.2), Inches(2.0), BG_CARD, color)
    txt(s, Inches(left + 0.15), Inches(1.9), Inches(1.9), Inches(0.3), title, 15, color, True, PP_ALIGN.CENTER)
    txt(s, Inches(left + 0.15), Inches(2.3), Inches(1.9), Inches(1.3), body, 12, LIGHT, False, PP_ALIGN.CENTER)

arrows = [(2.55, "SSO token"), (4.95, "channel token"), (7.35, "Bearer service_token"), (9.95, "Bearer downstream_token")]
for x, label in arrows:
    txt(s, Inches(x), Inches(2.55), Inches(0.5), Inches(0.25), "\u2192", 22, ACCENT, True, PP_ALIGN.CENTER)
    txt(s, Inches(x - 0.3), Inches(2.85), Inches(1.1), Inches(0.3), label, 9, DIM, False, PP_ALIGN.CENTER)

# Key principles
rect(s, Inches(0.3), Inches(4.2), Inches(12.7), Inches(3.0), BG_CARD, RGBColor(0x44, 0x44, 0x66))
txt(s, Inches(0.6), Inches(4.3), Inches(12), Inches(0.4), "Key Principles", 20, ACCENT, True)
principles = [
    ("1", "Token never leaves the OBO chain", "Each service receives a scoped assertion and exchanges it for the next hop."),
    ("2", "JWT validation lives in the Gateway", "Centralizes auth. Service trusts the gateway on an internal-only network."),
    ("3", "Service internal-only", "Not publicly accessible. Gateway is the sole ingress — the trust boundary."),
    ("4", "ContextVar carries the assertion", "Bound in middleware, available deep in the call stack for downstream OBO."),
]
y = Inches(4.8)
for num, title, desc in principles:
    rect(s, Inches(0.6), y, Inches(0.4), Inches(0.4), ACCENT)
    txt(s, Inches(0.6), y, Inches(0.4), Inches(0.4), num, 16, WHITE, True, PP_ALIGN.CENTER)
    txt(s, Inches(1.15), y, Inches(3.0), Inches(0.35), title, 14, WHITE, True)
    txt(s, Inches(4.2), y, Inches(8.5), Inches(0.5), desc, 12, LIGHT)
    y += Inches(0.62)

# ── SLIDE 4: Entra App Registrations ────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[BLANK])
set_bg(s, BG_DARK)
rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)
txt(s, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6), "Entra ID App Registrations", 32, WHITE, True)
txt(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.4),
    "Two app registrations + downstream resource awareness", 16, DIM)

# Service App
rect(s, Inches(0.5), Inches(1.6), Inches(5.5), Inches(2.8), BG_CARD, ACCENT)
txt(s, Inches(0.8), Inches(1.7), Inches(5.0), Inches(0.4), "Agentic Service App", 20, ACCENT, True)
for i, (k, v) in enumerate([
    ("Identifier URI", "api://<service-client-id>"),
    ("Sign-in audience", "AzureADMyOrg (single tenant)"),
    ("Token version", "requestedAccessTokenVersion: 2"),
    ("Exposed scope", "access_as_user (delegated)"),
    ("Required perms", "Downstream user_impersonation"),
    ("Client secret", "Yes (for MSAL ConfidentialClient OBO)"),
]):
    y = Inches(2.2) + i * Inches(0.35)
    txt(s, Inches(0.8), y, Inches(2.0), Inches(0.28), k, 11, DIM, True)
    txt(s, Inches(2.8), y, Inches(3.0), Inches(0.28), v, 11, LIGHT)

# Bot App
rect(s, Inches(6.5), Inches(1.6), Inches(6.3), Inches(2.8), BG_CARD, ACCENT2)
txt(s, Inches(6.8), Inches(1.7), Inches(5.8), Inches(0.4), "Gateway / Bot App", 20, ACCENT2, True)
for i, (k, v) in enumerate([
    ("Identifier URI", "api://botid-<bot-client-id>"),
    ("Sign-in audience", "AzureADMyOrg"),
    ("Exposed scope", "access_as_user (delegated)"),
    ("Required perms", "Service app access_as_user"),
    ("Redirect URI", "https://token.botframework.com/.auth/web/redirect"),
    ("Preauthorized", "Teams Desktop + Teams Web client IDs"),
]):
    y = Inches(2.2) + i * Inches(0.35)
    txt(s, Inches(6.8), y, Inches(2.2), Inches(0.28), k, 11, DIM, True)
    txt(s, Inches(9.0), y, Inches(3.6), Inches(0.28), v, 11, LIGHT)

# Admin consent
rect(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.5), BG_CARD, ACCENT3)
txt(s, Inches(0.8), Inches(4.8), Inches(11.5), Inches(0.4), "Admin Consent Required (both chains)", 18, ACCENT3, True)
txt(s, Inches(0.8), Inches(5.3), Inches(11.5), Inches(0.35),
    "Gateway/Bot app  \u2192  Service app access_as_user              \u2190  admin consent", 14, LIGHT)
txt(s, Inches(0.8), Inches(5.7), Inches(11.5), Inches(0.35),
    "Service app          \u2192  Downstream user_impersonation    \u2190  admin consent", 14, LIGHT)
txt(s, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.35),
    "Without admin consent: AADSTS65001 \u2014 The user or administrator has not consented", 12, RED)

# ── SLIDE 5: Gateway ───────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[BLANK])
set_bg(s, BG_DARK)
rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT2)
txt(s, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6), "Component 1: The M365 Gateway", 32, WHITE, True)
txt(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.4),
    "Stateless protocol adapter \u2014 bridges Bot Framework to your service's native API", 16, DIM)

# Key insight callout
rect(s, Inches(0.3), Inches(1.45), Inches(12.7), Inches(0.7), RGBColor(0x1E, 0x3A, 0x1E), ACCENT)
txt(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.55),
    "Key insight: You do NOT need to modify your existing service. The gateway absorbs JWT validation, "
    "protocol translation, and adaptation to your service's native API.", 13, ACCENT, True)

# Gateway flow
flow = [
    ("POST /api/messages", "Bot activity from\nCopilot/Teams", ACCENT2, 0.3),
    ("Route by type", "message \u2192 auth\ninvoke \u2192 SSO\nother \u2192 welcome", WHITE, 2.8),
    ("OBO #1 + JWT", "Token exchange\n+ validate JWT\n+ extract claims", ACCENT, 5.3),
    ("Translate &\nForward", "Map to service's\nnative API +\nX-User-* headers", ACCENT, 7.8),
    ("Extract reply", "Send activity\nback to Copilot", ACCENT2, 10.5),
]
for title, body, color, left in flow:
    rect(s, Inches(left), Inches(2.45), Inches(2.2), Inches(1.5), BG_CARD, color)
    txt(s, Inches(left + 0.1), Inches(2.5), Inches(2.0), Inches(0.4), title, 13, color, True, PP_ALIGN.CENTER)
    txt(s, Inches(left + 0.1), Inches(2.95), Inches(2.0), Inches(0.9), body, 11, LIGHT, False, PP_ALIGN.CENTER)

# Service client adapter concept
rect(s, Inches(0.3), Inches(4.2), Inches(6.3), Inches(3.0), BG_CARD, RGBColor(0x44, 0x44, 0x66))
txt(s, Inches(0.5), Inches(4.3), Inches(5.9), Inches(0.35), "Service Client = Adaptation Layer", 16, ACCENT2, True)
txt(s, Inches(0.5), Inches(4.7), Inches(5.9), Inches(0.5),
    "A small HTTP adapter class in the gateway maps Copilot conversations to your service:", 11, LIGHT)
for i, (cop, svc) in enumerate([
    ("conversation.id", "Session / thread / chat ID"),
    ("message text", "Request body (text, prompt, input, ...)"),
    ("reply extraction", "Response field (reply, output, content, ...)"),
    ("Bearer token", "Auth header, API key, or custom field"),
]):
    y = Inches(5.35) + i * Inches(0.36)
    txt(s, Inches(0.6), y, Inches(2.2), Inches(0.28), cop, 11, ACCENT2, True, PP_ALIGN.LEFT, "Consolas")
    txt(s, Inches(2.9), y, Inches(0.3), Inches(0.28), "\u2192", 12, DIM, False, PP_ALIGN.CENTER)
    txt(s, Inches(3.3), y, Inches(3.3), Inches(0.28), svc, 11, LIGHT)

# Implementation details
details = [
    ("JWT validation", "Verify signature, audience, issuer, expiry + extract claims", RED),
    ("Dual auth handlers", "Agentic channel + standard Teams connector paths", ACCENT2),
    ("Invoke handling", "No-op for SSO handshake or sign-in breaks", ACCENT3),
    ("Claim forwarding", "X-User-Id, X-User-UPN headers to service", ACCENT),
]
rect(s, Inches(6.9), Inches(4.2), Inches(6.1), Inches(3.0), BG_CARD, RGBColor(0x44, 0x44, 0x66))
txt(s, Inches(7.1), Inches(4.3), Inches(5.7), Inches(0.35), "Implementation Details", 16, ACCENT, True)
y = Inches(4.8)
for title, desc, color in details:
    txt(s, Inches(7.2), y, Inches(2.2), Inches(0.28), title, 12, color, True)
    txt(s, Inches(9.4), y, Inches(3.4), Inches(0.28), desc, 11, LIGHT)
    y += Inches(0.55)

# ── SLIDE 6: Service ─────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[BLANK])
set_bg(s, BG_DARK)
rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)
txt(s, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6), "Component 2: The Agentic Service", 32, WHITE, True)
txt(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.4),
    "Stateful API \u2014 knows nothing about Bot Framework or JWT validation", 16, DIM)

# Request flow
rect(s, Inches(0.3), Inches(1.5), Inches(6.2), Inches(5.5), BG_CARD, RGBColor(0x44, 0x44, 0x66))
txt(s, Inches(0.5), Inches(1.6), Inches(5.8), Inches(0.4), "Request Flow (Gateway Trust)", 18, ACCENT, True)
steps = [
    ("1", "Extract user identity", "From forwarded headers (X-User-Id, X-User-UPN)"),
    ("2", "Bind to ContextVar", "user_assertion + claims \u2192 async-safe ContextVar"),
    ("3", "Route to session", "owner_id from X-User-Id \u2192 403 on mismatch"),
    ("4", "Run agentic logic", "Any framework: MAF, LangChain, SK, custom"),
    ("5", "OBO #2 when needed", "ContextVar assertion \u2192 MSAL \u2192 downstream token"),
    ("6", "Return reply", '200 {session_id, reply, turns}'),
]
y = Inches(2.15)
for num, title, desc in steps:
    rect(s, Inches(0.5), y, Inches(0.35), Inches(0.35), ACCENT)
    txt(s, Inches(0.5), y, Inches(0.35), Inches(0.35), num, 13, WHITE, True, PP_ALIGN.CENTER)
    txt(s, Inches(1.0), y, Inches(2.3), Inches(0.3), title, 13, WHITE, True)
    txt(s, Inches(3.3), y, Inches(3.0), Inches(0.35), desc, 11, LIGHT)
    y += Inches(0.55)

# Minimum requirements (replaces JWT checklist)
rect(s, Inches(6.8), Inches(1.5), Inches(6.2), Inches(2.7), BG_CARD, RGBColor(0x44, 0x44, 0x66))
txt(s, Inches(7.0), Inches(1.6), Inches(5.8), Inches(0.4), "Minimum Service Requirements", 18, ACCENT2, True)
for i, item in enumerate([
    "Accept forwarded Bearer token (for OBO #2)",
    "Accept forwarded claim headers (X-User-Id, etc.)",
    "Session/thread identity across turns",
    "Message exchange endpoint (any shape)",
    "Internal-only ingress (not public)",
    "OBO #2 only if user-delegated access needed",
]):
    txt(s, Inches(7.2), Inches(2.1) + i * Inches(0.32), Inches(0.2), Inches(0.25), "\u2713", 11, ACCENT, True)
    txt(s, Inches(7.45), Inches(2.1) + i * Inches(0.32), Inches(5.3), Inches(0.25), item, 11, LIGHT)

# Framework freedom
rect(s, Inches(6.8), Inches(4.5), Inches(6.2), Inches(2.5), BG_CARD, RGBColor(0x44, 0x44, 0x66))
txt(s, Inches(7.0), Inches(4.6), Inches(5.8), Inches(0.4), "Your Agentic Logic Goes Here", 18, ACCENT, True)
for i, (fw, desc) in enumerate([
    ("Microsoft Agent Framework", "HandoffBuilder, ConcurrentBuilder, etc."),
    ("LangChain / LangGraph", "Chains, graphs, tool calling"),
    ("Semantic Kernel", "Planners and plugins"),
    ("Custom code", "Direct OpenAI SDK with tool loops"),
    ("Any other framework", "CrewAI, AutoGen, etc."),
]):
    txt(s, Inches(7.2), Inches(5.1) + i * Inches(0.32), Inches(2.8), Inches(0.25), fw, 12, ACCENT2, True)
    txt(s, Inches(10.0), Inches(5.1) + i * Inches(0.32), Inches(2.8), Inches(0.25), desc, 11, LIGHT)

# ── SLIDE 7: ContextVar + OBO ────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[BLANK])
set_bg(s, BG_DARK)
rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.06), PINK)
txt(s, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6), "ContextVar + OBO Pattern", 32, WHITE, True)
txt(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.4),
    "The core pattern that makes delegated auth work across any framework stack", 16, DIM)

# ContextVar code
rect(s, Inches(0.3), Inches(1.5), Inches(6.3), Inches(5.5), BG_CARD, RGBColor(0x44, 0x44, 0x66))
txt(s, Inches(0.5), Inches(1.6), Inches(5.9), Inches(0.4), "ContextVar: Bind + Reset", 18, PINK, True)
code1 = [
    "from contextvars import ContextVar, Token",
    "",
    "_USER_ASSERTION: ContextVar[str | None] =",
    "    ContextVar('user_assertion', default=None)",
    "",
    "def bind_identity(assertion, claims):",
    "    return (",
    "        _USER_ASSERTION.set(assertion),",
    "        _USER_CLAIMS.set(claims),",
    "    )",
    "",
    "# In middleware:",
    "a_tok, c_tok = bind_identity(raw_jwt, claims)",
    "try:",
    "    response = await call_next(request)",
    "finally:",
    "    _USER_ASSERTION.reset(a_tok)  # always!",
    "    _USER_CLAIMS.reset(c_tok)",
]
for i, line in enumerate(code1):
    c = ACCENT if line.startswith("#") else (ACCENT3 if "finally" in line or "always" in line else LIGHT)
    txt(s, Inches(0.6), Inches(2.1) + i * Inches(0.24), Inches(5.8), Inches(0.25), line, 11, c, False, PP_ALIGN.LEFT, "Consolas")

# OBO #2 code
rect(s, Inches(6.9), Inches(1.5), Inches(6.1), Inches(3.2), BG_CARD, RGBColor(0x44, 0x44, 0x66))
txt(s, Inches(7.1), Inches(1.6), Inches(5.7), Inches(0.4), "OBO #2: Downstream Token", 18, ACCENT3, True)
code2 = [
    "import msal",
    "",
    "app = msal.ConfidentialClientApplication(",
    "    client_id=SERVICE_CLIENT_ID,",
    "    client_credential=SERVICE_CLIENT_SECRET,",
    "    authority=f'https://login.microsoftonline",
    "              .com/{TENANT_ID}',",
    ")",
    "",
    "result = app.acquire_token_on_behalf_of(",
    "    user_assertion=_USER_ASSERTION.get(),",
    "    scopes=['<downstream>/.default'],",
    ")",
    "token = result['access_token']",
]
for i, line in enumerate(code2):
    txt(s, Inches(7.2), Inches(2.1) + i * Inches(0.24), Inches(5.6), Inches(0.25), line, 11, LIGHT, False, PP_ALIGN.LEFT, "Consolas")

# Why ContextVar
rect(s, Inches(6.9), Inches(5.0), Inches(6.1), Inches(2.0), BG_CARD, RGBColor(0x44, 0x44, 0x66))
txt(s, Inches(7.1), Inches(5.1), Inches(5.7), Inches(0.4), "Why ContextVar?", 16, ACCENT2, True)
for i, (t, d) in enumerate([
    ("Async-safe", "Each concurrent request gets its own scope"),
    ("Framework-agnostic", "Works with FastAPI, Flask, Django, raw asyncio"),
    ("Clean stack", "No need to thread assertion through every function"),
    ("Always reset", "finally block prevents leaking across requests"),
]):
    txt(s, Inches(7.2), Inches(5.55) + i * Inches(0.32), Inches(1.8), Inches(0.25), t, 11, ACCENT2, True)
    txt(s, Inches(9.0), Inches(5.55) + i * Inches(0.32), Inches(3.8), Inches(0.25), d, 11, LIGHT)

# ── SLIDE 8: Session + API Contract ──────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[BLANK])
set_bg(s, BG_DARK)
rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)
txt(s, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6), "Session Management & API Contract", 32, WHITE, True)

# Session mgmt
rect(s, Inches(0.3), Inches(1.2), Inches(6.3), Inches(5.8), BG_CARD, RGBColor(0x44, 0x44, 0x66))
txt(s, Inches(0.5), Inches(1.3), Inches(5.9), Inches(0.4), "Session Store (MVP: in-memory)", 18, ACCENT, True)
for i, (t, d) in enumerate([
    ("Owner isolation", "Every session has owner_id from X-User-Id.\nReject cross-user access with 403."),
    ("Turn windowing", "Cap stored turns at max_turns x 2.\nPrevents unbounded memory growth."),
    ("Concurrency lock", "asyncio.Lock() per session.\nPrevents interleaved turns."),
    ("Conversation mapping", "conversation.id maps 1:1 to session_id.\nAuto-create on 404."),
    ("Scaling path", "Replace with Redis / Cosmos DB.\nContextVar + OBO unaffected."),
]):
    y = Inches(1.85) + i * Inches(0.78)
    txt(s, Inches(0.5), y, Inches(2.2), Inches(0.3), t, 13, ACCENT2, True)
    txt(s, Inches(2.7), y, Inches(3.7), Inches(0.6), d, 11, LIGHT)

# API contract
rect(s, Inches(6.9), Inches(1.2), Inches(6.1), Inches(5.8), BG_CARD, RGBColor(0x44, 0x44, 0x66))
txt(s, Inches(7.1), Inches(1.3), Inches(5.7), Inches(0.4), "Service API Contract", 18, ACCENT, True)
txt(s, Inches(7.1), Inches(1.75), Inches(5.7), Inches(0.3), "Reference contract (adapt the gateway if yours differs):", 13, DIM)

for i, (method, path, desc, ret) in enumerate([
    ("POST", "/api/chat/sessions", "Create session", "Returns session_id"),
    ("POST", "/api/chat/sessions/{id}/messages", "Send message", "Returns reply + turns"),
    ("GET", "/api/chat/sessions/{id}", "Get session", "Returns turn history"),
    ("GET", "/healthz", "Health check", "No auth required"),
]):
    y = Inches(2.3) + i * Inches(0.62)
    mc = ACCENT if method == "POST" else ACCENT2
    txt(s, Inches(7.2), y, Inches(0.6), Inches(0.28), method, 11, mc, True)
    txt(s, Inches(7.85), y, Inches(3.2), Inches(0.28), path, 11, LIGHT, False, PP_ALIGN.LEFT, "Consolas")
    txt(s, Inches(7.2), y + Inches(0.28), Inches(5.5), Inches(0.25), f"{desc} \u2014 {ret}", 10, DIM)

txt(s, Inches(7.1), Inches(4.95), Inches(5.7), Inches(0.35), "Response Codes", 14, ACCENT3, True)
for i, (code, desc, c) in enumerate([
    ("200", "Success with reply", ACCENT),
    ("401", "Bad or missing token", RED),
    ("403", "Session belongs to different user", RED),
    ("404", "Session not found (gateway auto-creates + retries)", ACCENT3),
]):
    txt(s, Inches(7.2), Inches(5.4) + i * Inches(0.32), Inches(0.5), Inches(0.25), code, 12, c, True, PP_ALIGN.LEFT, "Consolas")
    txt(s, Inches(7.8), Inches(5.4) + i * Inches(0.32), Inches(5.0), Inches(0.25), desc, 11, LIGHT)

# ── SLIDE 9: M365 App Package ────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[BLANK])
set_bg(s, BG_DARK)
rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.06), PURPLE)
txt(s, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6), "Component 3: M365 App Package", 32, WHITE, True)
txt(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.4),
    "ZIP containing manifest + icons that registers the gateway as a Custom Engine Agent", 16, DIM)

# Package contents
rect(s, Inches(0.3), Inches(1.5), Inches(4.0), Inches(2.3), BG_CARD, PURPLE)
txt(s, Inches(0.5), Inches(1.6), Inches(3.6), Inches(0.4), "Package Contents (ZIP)", 16, PURPLE, True)
for i, (f, d) in enumerate([
    ("manifest.json", "Bot + SSO + CEA config"),
    ("color.png", "192x192 color icon"),
    ("outline.png", "32x32 outline icon"),
]):
    txt(s, Inches(0.6), Inches(2.15) + i * Inches(0.35), Inches(1.5), Inches(0.25), f, 12, WHITE, True, PP_ALIGN.LEFT, "Consolas")
    txt(s, Inches(2.1), Inches(2.15) + i * Inches(0.35), Inches(2.0), Inches(0.25), d, 11, LIGHT)

# Critical manifest fields
rect(s, Inches(4.8), Inches(1.5), Inches(8.2), Inches(5.5), BG_CARD, RGBColor(0x44, 0x44, 0x66))
txt(s, Inches(5.0), Inches(1.6), Inches(7.8), Inches(0.4), "Critical Manifest Fields", 18, PURPLE, True)
for i, (f, v, d) in enumerate([
    ("bots[0].botId", "BOT_APP_ID", "Bot registration app ID"),
    ("bots[0].scopes", '["personal"]', "Personal scope for Copilot"),
    ("webApplicationInfo.id", "BOT_SSO_APP_ID", "SSO app ID (same as bot)"),
    ("webApplicationInfo.resource", "api://botid-<id>", "SSO resource URI"),
    ("copilotAgents...type", "bot", "Custom Engine Agent type"),
    ("copilotAgents...functionsAs", "agentOnly", "Agent-only mode"),
    ("copilotAgents...id", "BOT_APP_ID", "Same bot app ID"),
    ("validDomains[0]", "<gateway>.azurecontainerapps.io", "Gateway hostname"),
    ("validDomains[1]", "token.botframework.com", "Bot Framework token domain"),
]):
    y = Inches(2.15) + i * Inches(0.35)
    txt(s, Inches(5.1), y, Inches(2.7), Inches(0.28), f, 10, ACCENT2, False, PP_ALIGN.LEFT, "Consolas")
    txt(s, Inches(7.8), y, Inches(2.7), Inches(0.28), v, 10, LIGHT, False, PP_ALIGN.LEFT, "Consolas")
    txt(s, Inches(10.5), y, Inches(2.3), Inches(0.28), d, 10, DIM)

# Bot OAuth note
rect(s, Inches(0.3), Inches(4.1), Inches(4.0), Inches(2.9), BG_CARD, ACCENT3)
txt(s, Inches(0.5), Inches(4.2), Inches(3.6), Inches(0.35), "Azure Bot OAuth Connection", 14, ACCENT3, True)
for i, line in enumerate([
    "Service: Aadv2",
    "Client ID: BOT_APP_ID",
    "Client Secret: BOT_APP_PASSWORD",
    "Scope: <service>/access_as_user",
    "  offline_access openid profile",
    "TokenExchangeUrl:",
    "  api://botid-<BOT_APP_ID>",
]):
    c = ACCENT2 if line.startswith("  ") else LIGHT
    txt(s, Inches(0.5), Inches(4.65) + i * Inches(0.28), Inches(3.6), Inches(0.25), line, 11, c)

# ── SLIDE 10: Deployment Sequence ────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[BLANK])
set_bg(s, BG_DARK)
rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)
txt(s, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6), "Deployment Sequence", 32, WHITE, True)
txt(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.4),
    "Follow this exact order \u2014 each step depends on previous outputs", 16, DIM)

deploy = [
    ("1", "Entra app registrations", "e1f5fe", "Identity"),
    ("2", "Grant admin consent", "fff3e0", "Identity"),
    ("3", "Azure Bot + OAuth", "e1f5fe", "Identity"),
    ("4", "Prepare downstream", "e8f5e9", "Infra"),
    ("5", "Deploy Service (ACA)", "e8f5e9", "Deploy"),
    ("6", "Validate service", "c8e6c9", "Test"),
    ("7", "Deploy Gateway (ACA)", "e8f5e9", "Deploy"),
    ("8", "Validate gateway", "c8e6c9", "Test"),
    ("9", "Build app package", "f3e5f5", "Package"),
    ("10", "Upload to catalog", "f3e5f5", "Package"),
    ("11", "Wire Bot endpoint", "e1f5fe", "Wire"),
    ("12", "Test in Copilot", "c8e6c9", "Test"),
]
for i, (num, label, hx, cat) in enumerate(deploy):
    row, col = divmod(i, 4)
    x = Inches(0.5) + col * Inches(3.15)
    y = Inches(1.5) + row * Inches(1.65)
    r, g, b = int(hx[0:2], 16) // 3, int(hx[2:4], 16) // 3, int(hx[4:6], 16) // 3
    border = RGBColor(int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16))
    rect(s, x, y, Inches(2.9), Inches(1.3), RGBColor(r, g, b), border)
    rect(s, x, y, Inches(0.4), Inches(1.3), border)
    txt(s, x, y + Inches(0.3), Inches(0.4), Inches(0.4), num, 16, BG_DARK, True, PP_ALIGN.CENTER)
    txt(s, x + Inches(0.5), y + Inches(0.15), Inches(2.3), Inches(0.3), label, 14, WHITE, True)
    txt(s, x + Inches(0.5), y + Inches(0.5), Inches(2.3), Inches(0.25), cat, 10, DIM)

txt(s, Inches(0.5), Inches(6.6), Inches(12), Inches(0.4),
    "Each step outputs values (app IDs, secrets, URLs) consumed by subsequent steps. Store in .env.", 13, DIM)

# ── SLIDE 11: Security Checklist ─────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[BLANK])
set_bg(s, BG_DARK)
rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.06), RED)
txt(s, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6), "Security Checklist", 32, WHITE, True)
txt(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.4), "Defense in depth \u2014 every layer validates", 16, DIM)

checks = [
    ("JWT validated in Gateway", "Audience, issuer, signature, expiry", ACCENT),
    ("JWT issuer validated", "Accept only tenant STS URIs", ACCENT),
    ("JWT signature verified", "RS256 via JWKS, not hardcoded", ACCENT),
    ("Service internal-only", "Not publicly accessible, gateway sole ingress", RED),
    ("ContextVar reset in finally", "Prevents assertion leaking", PINK),
    ("Session owner isolation", "X-User-Id based, 403 on mismatch", ACCENT2),
    ("Downstream data scoped", "Parameterized queries, RLS, no raw SQL", ACCENT3),
    ("Secrets in ACA secrets", "secretref:, not plain env vars", ACCENT3),
    ("Invoke handling", "Clean return for SSO handshake", ACCENT2),
    ("Admin consent granted", "Both OBO chains consented", RED),
]
for i, (title, desc, color) in enumerate(checks):
    col, row = divmod(i, 5)
    x = Inches(0.3) + col * Inches(6.5)
    y = Inches(1.5) + row * Inches(0.95)
    rect(s, x, y, Inches(6.2), Inches(0.8), BG_CARD, RGBColor(0x44, 0x44, 0x66))
    rect(s, x, y, Inches(0.08), Inches(0.8), color)
    txt(s, x + Inches(0.25), y + Inches(0.05), Inches(5.7), Inches(0.3), title, 14, color, True)
    txt(s, x + Inches(0.25), y + Inches(0.4), Inches(5.7), Inches(0.3), desc, 11, LIGHT)

# ── SLIDE 12: Troubleshooting ────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[BLANK])
set_bg(s, BG_DARK)
rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT3)
txt(s, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6), "Common Pitfalls & Troubleshooting", 32, WHITE, True)

for i, (sym, cause, fix) in enumerate([
    ("AADSTS65001: not consented", "Missing admin consent", "Grant admin consent for both OBO chains"),
    ("Token is not exchangeable", "Bot OAuth misconfigured", "Recreate OAuth connection; new conversation"),
    ("OBO #2 invalid_grant", "Assertion expired / aud mismatch", "Ensure tokenVersion: 2 + correct audience"),
    ("Gateway returns 401", "JWT validation rejects OBO #1 token", "Check SERVICE_EXPECTED_AUDIENCE, issuer"),
    ("Session 403 on 2nd turn", "owner_id mismatch between turns", "Verify X-User-Id header is consistent"),
    ("/healthz returns 401", "Health not excluded from auth", "Add path exclusion for /healthz"),
    ("Invoke shows error", "Missing invoke handler", "Add no-op route with is_invoke=True"),
]):
    y = Inches(1.2) + i * Inches(0.82)
    rect(s, Inches(0.3), y, Inches(12.7), Inches(0.72), BG_CARD, RGBColor(0x44, 0x44, 0x66))
    txt(s, Inches(0.5), y + Inches(0.05), Inches(4.0), Inches(0.3), sym, 11, RED, True, PP_ALIGN.LEFT, "Consolas")
    txt(s, Inches(4.7), y + Inches(0.05), Inches(3.5), Inches(0.3), cause, 11, ACCENT3, True)
    txt(s, Inches(4.7), y + Inches(0.38), Inches(8.0), Inches(0.3), fix, 11, ACCENT)

# ── SLIDE 13: Extending ──────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[BLANK])
set_bg(s, BG_DARK)
rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)
txt(s, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6), "Extending the Pattern", 32, WHITE, True)

extensions = [
    ("Fronting an Existing Service", PINK, [
        "Do NOT modify the existing service",
        "Gateway handles JWT validation for you",
        "Write a service client adapter in gateway",
        "Deploy service on internal-only ingress",
    ]),
    ("Adding a New Downstream", ACCENT2, [
        "Add delegated permission to service app",
        "Grant admin consent for the new resource",
        "Add acquire_X_token() using same ContextVar",
        "Call from agentic logic when needed",
    ]),
    ("Scaling Beyond Single-Replica", ACCENT, [
        "Replace InMemorySessionStore with Redis/Cosmos",
        "ContextVar + OBO is per-request, not per-replica",
        "Gateway already stateless \u2014 scale horizontally",
        "Session affinity via conversation_id",
    ]),
    ("Adding a Second Agent", ACCENT3, [
        "Deploy second gateway + service pair",
        "Own Entra app registrations + Bot resource",
        "Separate or multi-bot M365 app package",
        "Same OBO chain pattern applies independently",
    ]),
]
for i, (title, color, items) in enumerate(extensions):
    x = Inches(0.3) + i * Inches(3.2)
    rect(s, x, Inches(1.2), Inches(3.0), Inches(5.5), BG_CARD, color)
    txt(s, x + Inches(0.15), Inches(1.3), Inches(2.7), Inches(0.4), title, 14, color, True)
    for j, item in enumerate(items):
        y = Inches(1.95) + j * Inches(0.7)
        rect(s, x + Inches(0.15), y, Inches(0.3), Inches(0.3), color)
        txt(s, x + Inches(0.15), y, Inches(0.3), Inches(0.3), str(j + 1), 11, BG_DARK, True, PP_ALIGN.CENTER)
        txt(s, x + Inches(0.55), y, Inches(2.3), Inches(0.6), item, 11, LIGHT)

# ── Save ─────────────────────────────────────────────────────────────
out = "docs/m365-agentic-service-developer-guide.pptx"
prs.save(out)
print(f"Saved: {out}")
